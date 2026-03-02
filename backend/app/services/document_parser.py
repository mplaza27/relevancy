from __future__ import annotations

from pathlib import Path


def extract_pdf(file_path: Path) -> str:
    """Extract text from a PDF using pymupdf4llm (outputs Markdown)."""
    import pymupdf4llm

    return pymupdf4llm.to_markdown(str(file_path))


def extract_pptx(file_path: Path) -> str:
    """Extract text from a PPTX file, slide by slide."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(file_path))
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)


def extract_docx(file_path: Path) -> str:
    """Extract text from a DOCX file preserving paragraph and table order."""
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(file_path))
    parts: list[str] = []

    for child in doc.element.body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            # Paragraph
            from docx.text.paragraph import Paragraph
            para = Paragraph(child, doc)
            text = para.text.strip()
            if text:
                parts.append(text)

        elif tag == "tbl":
            # Table
            from docx.table import Table
            table = Table(child, doc)
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".txt": lambda p: p.read_text(encoding="utf-8", errors="replace"),
    ".md": lambda p: p.read_text(encoding="utf-8", errors="replace"),
}

SUPPORTED_EXTENSIONS = frozenset(_EXTRACTORS.keys())


def extract_text(file_path: Path) -> str:
    """Route a file to the appropriate extractor based on its extension.

    Raises ValueError for unsupported file types.
    """
    suffix = file_path.suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(
            f"Unsupported file type: {suffix!r}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )
    return extractor(file_path)
