"""Tests for document_parser.py.

PDF, PPTX, DOCX tests create minimal in-memory fixtures using the respective libraries.
"""

import io
import tempfile
from pathlib import Path

import pytest

from app.services.document_parser import extract_text


class TestExtractText:
    def test_unsupported_extension_raises(self):
        with pytest.raises(ValueError, match="Unsupported"):
            extract_text(Path("file.xyz"))

    def test_txt_extraction(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("Hello, world!", encoding="utf-8")
        result = extract_text(f)
        assert result == "Hello, world!"

    def test_md_extraction(self, tmp_path):
        f = tmp_path / "test.md"
        f.write_text("# Title\n\nContent here.", encoding="utf-8")
        result = extract_text(f)
        assert "Title" in result
        assert "Content" in result


class TestExtractDocx:
    def test_basic_docx(self, tmp_path):
        pytest.importorskip("docx")
        from docx import Document

        doc = Document()
        doc.add_paragraph("First paragraph about cardiology.")
        doc.add_paragraph("Second paragraph about anatomy.")
        docx_path = tmp_path / "test.docx"
        doc.save(str(docx_path))

        result = extract_text(docx_path)
        assert "First paragraph" in result
        assert "Second paragraph" in result

    def test_docx_with_table(self, tmp_path):
        pytest.importorskip("docx")
        from docx import Document

        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Cell A"
        table.cell(1, 1).text = "Cell B"
        docx_path = tmp_path / "table.docx"
        doc.save(str(docx_path))

        result = extract_text(docx_path)
        assert "Header 1" in result
        assert "Cell A" in result


class TestExtractPptx:
    def test_basic_pptx(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(slide_layout)

        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.text = "Cardiology lecture content here."

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        result = extract_text(pptx_path)
        assert "Slide 1" in result
        assert "Cardiology" in result
